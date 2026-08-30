# -*- coding: utf-8 -*-
"""Rapport PowerPoint PANORAMA - Defi 2 (10 slides, palette drapeau togolais)."""
import os, sys, json
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = r"C:\Projet DATA IA\TOGO-ENVIRONNEMENT 2"
DATA = os.path.join(BASE, "data")
IMG = os.path.join(BASE, "rapport", "img")
ICONS = os.path.join(BASE, "assets", "icons")
LOGO = os.path.join(BASE, "assets", "logo.jpg")
K = json.load(open(os.path.join(DATA, "kpis.json"), encoding="utf-8"))

# palette drapeau
GREEN = RGBColor(0x0A, 0x7C, 0x46); GREEN_D = RGBColor(0x07, 0x5E, 0x35)
YELLOW = RGBColor(0xFF, 0xCE, 0x00); RED = RGBColor(0xD2, 0x10, 0x34)
GOLD = RGBColor(0xE0, 0xA4, 0x22); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x25, 0x1A); MUT = RGBColor(0x5E, 0x68, 0x5C)
PAPER = RGBColor(0xF3, 0xF5, 0xEF); LINE = RGBColor(0xDD, 0xE1, 0xD7)
GREENSOFT = RGBColor(0xE4, 0xF1, 0xEA); REDSOFT = RGBColor(0xFB, 0xE4, 0xE7)
GOLDSOFT = RGBColor(0xFB, 0xF1, 0xD9)
HEAD = "Cambria"; BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=WHITE):
    s = prs.slides.add_slide(BL)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s


def noline(shp):
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def rect(s, l, t, w, h, fill, radius=0.09, line=None, lw=1):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try: shp.adjustments[0] = radius
    except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.02):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs[0], tuple): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        p.line_spacing = line_sp
        for (t_, sz, bold, col, *rest) in para:
            r = p.add_run(); r.text = t_
            r.font.name = rest[1] if len(rest) > 1 else BODY
            r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col
            r.font.italic = rest[0] if rest else False
    return tb


def icon_circle(s, name, cx, cy, d, tint):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = tint; o.line.fill.background(); o.shadow.inherit = False
    isz = d * 0.56
    p = os.path.join(ICONS, name)
    if os.path.exists(p):
        s.shapes.add_picture(p, Inches(cx - isz / 2), Inches(cy - isz / 2), Inches(isz), Inches(isz))


def pic(s, path, l, t, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(l), Inches(t), **kw)


# ============================================================ S1 COVER
s = slide(GREEN)
rect(s, 0.7, 0.62, 1.15, 1.15, WHITE, radius=0.18)
pic(s, LOGO, 0.78, 0.7, w=0.99, h=0.99)
txt(s, 2.15, 0.9, 9, 0.4, [[("DATA CHALLENGE ENVIRONNEMENT  ·  DÉFI 2", 13, True, YELLOW)]])
txt(s, 2.13, 1.28, 10, 0.5, [[("Développé par Bastou OURO-TAGBA", 12, False, RGBColor(0xCF,0xE6,0xD8))]])
txt(s, 0.7, 2.7, 12, 1.5, [[("PANORAMA", 76, True, WHITE, False, HEAD)]])
txt(s, 0.72, 4.05, 11.5, 0.8, [[("Énergie ", 30, True, WHITE), ("&", 30, True, YELLOW),
                                (" transition écologique au Togo", 30, True, WHITE)]])
txt(s, 0.74, 4.85, 11, 0.6, [[("Électrifier sans déforester", 20, False, RGBColor(0xDCEAE0>>16 & 255,0,0) if False else RGBColor(0xDC,0xEA,0xE0), True)]])
txt(s, 0.72, 6.6, 12, 0.4, [[("Togo AI Lab / MESPTN  ·  Données ouvertes opendata.gouv.tg (WDI, Banque mondiale)", 12, False, RGBColor(0xB8,0xD8,0xC6))]])

# ============================================================ S2 CONTEXTE
s = slide(WHITE)
txt(s, 0.7, 0.55, 8, 0.7, [[("Le contexte", 34, True, INK, False, HEAD)]])
txt(s, 0.7, 1.25, 7.1, 0.4, [[("Électrifier tout le pays d'ici 2030, sans sacrifier la forêt", 15, False, MUT)]])
paras = [
    [("L'ambition.  ", 14, True, GREEN), ("Le Togo vise l'accès universel à l'électricité en 2030, "
     "tout en développant les énergies propres et en protégeant son environnement.", 14, False, INK)],
    [("Le contraste.  ", 14, True, RED), ("Les villes sont bien électrifiées, les campagnes très en retard. "
     "Et l'immense majorité des ménages cuisine encore au bois et au charbon, ce qui pèse sur des forêts "
     "déjà fragilisées par la hausse des températures.", 14, False, INK)],
    [("La question.  ", 14, True, INK), ("Comment apporter l'électricité dans les villages, développer les "
     "énergies propres et protéger les forêts, à partir des données ?", 14, False, INK)],
]
txt(s, 0.7, 2.0, 7.15, 4.6, paras, sp_after=12, line_sp=1.08)
# carte des 6 sources (droite)
rect(s, 8.25, 1.7, 4.4, 5.2, PAPER, radius=0.06)
txt(s, 8.6, 2.0, 3.8, 0.4, [[("6 JEUX DE DONNÉES OUVERTS", 12, True, GREEN)]])
sources = [("bolt_good.png", "Accès électricité & cuisson, population, économie"),
           ("factory_warn.png", "Émissions par secteur (gaz à effet de serre)"),
           ("thermo_warn.png", "Températures mensuelles · 10 villes"),
           ("tree_good.png", "Énergies renouvelables & biomasse"),
           ("chartup_warn.png", "CO₂ de l'électricité (1970-2022)"),
           ("leaf_good.png", "53 zones protégées & forêts classées")]
y = 2.55
for ic, lab in sources:
    icon_circle(s, ic, 9.0, y + 0.22, 0.62, GREENSOFT)
    txt(s, 9.45, y, 3.05, 0.7, [[(lab, 12.5, False, INK)]], anchor=MSO_ANCHOR.MIDDLE, line_sp=1.0)
    y += 0.72

# ============================================================ S3 DIAGNOSTIC
s = slide(PAPER)
txt(s, 0.7, 0.55, 11, 0.7, [[("Le diagnostic en un écran", 34, True, INK, False, HEAD)]])
txt(s, 0.7, 1.25, 11, 0.4, [[("Quatre chiffres qui résument l'urgence, et fixent le cap.", 15, False, MUT)]])
stats = [
    ("3,86 M", RED, "bolt_crit.png", REDSOFT, "Togolais sans électricité", "96 % de ce déficit est rural (2022)"),
    ("88 %", RED, "flame_crit.png", REDSOFT, "cuisinent sans énergie propre", "soit 7,8 M de personnes au bois / charbon"),
    ("−1 554", GOLD, "tree_warn.png", GOLDSOFT, "km² de forêts perdus depuis 1990", "−11 % du couvert · ~50 km²/an"),
    ("94 %", GOLD, "power_warn.png", GOLDSOFT, "des entreprises subissent des coupures", "~5,5 coupures/mois (2016)"),
]
cw, gap, x0, y0, ch = 2.86, 0.18, 0.7, 2.05, 3.0
for i, (num, col, ic, tint, lab, sub) in enumerate(stats):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, WHITE, radius=0.08)
    icon_circle(s, ic, x + 0.55, y0 + 0.6, 0.72, tint)
    txt(s, x + 0.28, y0 + 1.05, cw - 0.5, 0.9, [[(num, 40, True, col, False, HEAD)]])
    txt(s, x + 0.28, y0 + 1.95, cw - 0.5, 0.7, [[(lab, 13.5, True, INK)]], line_sp=1.0)
    txt(s, x + 0.28, y0 + 2.55, cw - 0.5, 0.4, [[(sub, 10.5, False, MUT)]], line_sp=1.0)
rect(s, 0.7, 5.55, 11.93, 1.15, GREEN, radius=0.07)
txt(s, 1.05, 5.72, 11.3, 0.85, [[("La thèse.  ", 16, True, YELLOW), ("Électrifier sans déforester : "
    "l'off-grid solaire pour l'accès, la cuisson propre pour épargner la forêt, un mix solaire pour ne pas "
    "carboniser la transition.", 16, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE, line_sp=1.05)


# ---- gabarit slide analyse : titre + graphe gauche + insights droite ----
def analyse(title, tag, chart, kicker, kcol, points, reco):
    s = slide(WHITE)
    txt(s, 0.7, 0.5, 11.9, 0.7, [[(title, 30, True, INK, False, HEAD)]])
    txt(s, 0.72, 1.15, 11.9, 0.4, [[(tag, 14, False, MUT)]])
    rect(s, 0.7, 1.75, 7.15, 5.15, PAPER, radius=0.05)
    pic(s, os.path.join(IMG, chart), 0.95, 2.05, w=6.65)
    txt(s, 8.2, 1.85, 4.5, 0.4, [[(kicker, 12, True, kcol)]])
    yy = 2.35
    for head, body in points:
        txt(s, 8.2, yy, 4.5, 0.4, [[(head, 15, True, INK)]])
        tb = txt(s, 8.2, yy + 0.38, 4.5, 1.0, [[(body, 13, False, INK)]], line_sp=1.05)
        yy += 0.42 + 0.30 * (1 + len(body) // 46)
    rect(s, 8.2, 6.02, 4.45, 0.95, GREENSOFT, radius=0.1)
    txt(s, 8.45, 6.14, 4.0, 0.75, [[("▸ ", 13, True, GREEN), (reco, 12.5, False, INK)]],
        anchor=MSO_ANCHOR.MIDDLE, line_sp=1.02)
    return s


# ============================================================ S4 ELEC
analyse("La fracture de la lumière", "Ville et campagne face à l'accès, et un réseau qui flanche.",
        "fracture.png", "ÉLECTRIFICATION & RÉSEAU", RED,
        [("Un écart qui se creuse", "72 points de % séparent ville (96 %) et campagne (25 %) : un écart ×2,2 depuis 2000."),
         ("Une projection qui alerte", "Au rythme 2010-2022, le rural n'atteint l'accès universel que vers ~2089."),
         ("Un réseau peu fiable", "94 % des entreprises subissent des coupures (~5,5/mois).")],
        "Solaire hors-réseau : mini-centrales et kits, la seule voie pour tenir le calendrier.")

# ============================================================ S5 FORETS
analyse("Le feu de bois & le recul des forêts", "La cuisson au bois, ou l'absence de choix qui pèse sur la forêt.",
        "deforestation.png", "CUISSON & FORÊTS", GOLD,
        [("Presque aucune cuisson propre", "En campagne, 0,9 % seulement des ménages cuisinent proprement."),
         ("Une énergie encore très bois", "La biomasse pèse 44 % de l'énergie du pays (2014)."),
         ("Une forêt qui recule", "−1 554 km² depuis 1990 (−11 %). Chaque repas prélève sur le couvert.")],
        "Diffuser GPL et foyers améliorés : alléger directement la pression sur les forêts.")

# ============================================================ S6 EMISSIONS
analyse("L'empreinte carbone : la nuance", "D'où viennent vraiment les gaz à effet de serre du Togo ?",
        "emissions.png", "ÉMISSIONS", GREEN,
        [("L'énergie pèse peu", "Elle ne représente que 6 % des émissions ; l'agriculture et les terres dominent (88 %)."),
         ("Un levier d'équité", "Électrifier le rural est d'abord une question d'équité et de protection des forêts."),
         ("Mais rester propre", "Le CO₂ de l'électricité augmente avec le raccordement.")],
        "Garder une électricité solaire : ne pas troquer le bois contre le fossile.")

# ============================================================ S7 CLIMAT
analyse("Du Sud au Nord", "Où il fait chaud, et comment la température évolue (10 villes).",
        "climat.png", "CLIMAT", GOLD,
        [("Un gradient net", "Jusqu'à +7,2 °C d'écart entre villes : il fait bien plus chaud au Nord."),
         ("Un réchauffement d'ensemble", "Sur 2013-2019, 8 villes sur 10 se réchauffent (+0,45 °C/décennie en moyenne)."),
         ("Le lien énergie", "Le Nord, le plus chaud, est aussi le moins électrifié, plus de besoins de refroidissement.")],
        "Renforce l'urgence d'électrifier proprement les campagnes du Nord.")

# ============================================================ S8 RECOS
s = slide(GREEN)
txt(s, 0.7, 0.55, 11, 0.7, [[("Agir, sobrement", 34, True, WHITE, False, HEAD)]])
txt(s, 0.72, 1.25, 11, 0.4, [[("Quatre leviers pour électrifier sans déforester.", 15, False, RGBColor(0xCF,0xE6,0xD8))]])
recos = [
    ("solar_lv.png", "01  Solariser les villages", "Mini-centrales et kits solaires là où le réseau est lent, coûteux et peu fiable.", "Savanes, Kara, Centrale"),
    ("pot_lv.png", "02  Sortir du bois de cuisson", "Diffuser GPL et foyers améliorés : moins de pression sur les forêts et la santé.", "Ménages ruraux, national"),
    ("cleanpower_lv.png", "03  Garder une électricité propre", "Privilégier le solaire pour ne pas remplacer le bois par le fossile.", "Mix électrique national"),
    ("shield_lv.png", "04  Protéger les forêts sous pression", "Cibler les Plateaux et la Centrale, qui portent l'essentiel du couvert classé.", "Plateaux & Centrale"),
]
cw, ch, gx, gy, x0, y0 = 5.85, 2.35, 0.22, 0.22, 0.7, 2.0
for i, (ic, head, body, cible) in enumerate(recos):
    x = x0 + (i % 2) * (cw + gx); y = y0 + (i // 2) * (ch + gy)
    rect(s, x, y, cw, ch, WHITE, radius=0.06)
    icon_circle(s, ic, x + 0.72, y + 0.72, 0.9, GREENSOFT)
    txt(s, x + 1.35, y + 0.28, cw - 1.6, 0.5, [[(head, 16.5, True, INK, False, HEAD)]])
    txt(s, x + 1.35, y + 0.86, cw - 1.6, 1.0, [[(body, 12.5, False, INK)]], line_sp=1.05)
    txt(s, x + 1.35, y + 1.9, cw - 1.6, 0.35, [[("CIBLE : ", 10, True, GREEN), (cible, 10, False, MUT)]])

# ============================================================ S9 DASHBOARD
s = slide(WHITE)
txt(s, 0.7, 0.55, 11.9, 0.7, [[("Le tableau de bord PANORAMA", 32, True, INK, False, HEAD)]])
txt(s, 0.72, 1.2, 11.9, 0.4, [[("Une salle de contrôle interactive", 15, False, MUT)]])
secs = [("bolt_good.png", "Vue", "Le Togo en un écran"), ("bolt_crit.png", "Élec", "Fracture & réseau"),
        ("tree_good.png", "Bois", "Cuisson & forêts"), ("factory_warn.png", "GES", "Émissions"),
        ("thermo_warn.png", "Climat", "Sud → Nord"), ("shield_lv.png", "Agir", "Recommandations")]
cw, gx, x0, y0 = 3.87, 0.16, 0.7, 2.0
for i, (ic, t_, d_) in enumerate(secs):
    x = x0 + (i % 3) * (cw + gx); y = y0 + (i // 3) * 1.15
    rect(s, x, y, cw, 1.0, PAPER, radius=0.1)
    icon_circle(s, ic, x + 0.6, y + 0.5, 0.66, GREENSOFT)
    txt(s, x + 1.05, y + 0.16, cw - 1.2, 0.4, [[(t_, 15, True, INK)]])
    txt(s, x + 1.05, y + 0.55, cw - 1.2, 0.35, [[(d_, 11.5, False, MUT)]])
rect(s, 0.7, 4.5, 11.93, 2.4, GREENSOFT, radius=0.05)
txt(s, 1.0, 4.75, 5.6, 0.4, [[("CE QUI LE REND VIVANT", 12, True, GREEN)]])
feats = [[("Carte pilote.  ", 13.5, True, INK), ("53 forêts classées + 10 stations ; cliquer une forêt filtre tout.", 13.5, False, INK)],
         [("Filtre régional.  ", 13.5, True, INK), ("Les KPI, la carte et l'inventaire se recalculent par région.", 13.5, False, INK)],
         [("Honnêteté des données.  ", 13.5, True, INK), ("Chaque vue indique quand une donnée n'existe qu'au niveau national.", 13.5, False, INK)]]
txt(s, 1.0, 5.2, 6.4, 1.6, feats, sp_after=8, line_sp=1.05)
rect(s, 7.7, 4.9, 4.6, 1.65, WHITE, radius=0.08, line=LINE, lw=1)
txt(s, 8.0, 5.12, 4.0, 0.4, [[("ACCÈS", 12, True, GREEN)]])
txt(s, 8.0, 5.5, 4.05, 1.0, [[("Application web (Python / Dash).", 12.5, False, INK)],
    [("pip install -r requirements.txt", 12, False, INK, False, "Consolas")],
    [("python app.py  →  localhost:8050", 12, False, INK, False, "Consolas")]], sp_after=4, line_sp=1.05)

# ============================================================ S10 METHODO
s = slide(GREEN_D)
txt(s, 0.7, 0.55, 11.9, 0.7, [[("Méthode, sources & limites", 32, True, WHITE, False, HEAD)]])
txt(s, 0.72, 1.2, 11.9, 0.4, [[("Transparence sur les données, et sur ce qu'elles ne disent pas.", 15, False, RGBColor(0xB8,0xD8,0xC6))]])
cols = [
    ("SOURCES", ["6 jeux de données ouverts du défi (WDI / opendata.gouv.tg).",
                 "1 seul ajout externe, purement cartographique : frontières des 5 régions (geoBoundaries)."]),
    ("MÉTHODE", ["Nettoyage et calculs reproductibles (scripts fournis).",
                 "KPI et tendances calculés depuis les données, jamais saisis à la main."]),
    ("LIMITES ASSUMÉES", ["Accès électricité, cuisson et émissions : données nationales, non déclinées par région.",
                          "Températures : 10 villes, série courte 2013-2019."]),
]
x0, cw, gx = 0.7, 3.87, 0.16
for i, (head, items) in enumerate(cols):
    x = x0 + i * (cw + gx)
    rect(s, x, 2.0, cw, 3.7, RGBColor(0x0C, 0x6B, 0x3E), radius=0.06)
    txt(s, x + 0.3, 2.25, cw - 0.6, 0.4, [[(head, 13, True, YELLOW)]])
    paras = []
    for it in items:
        paras.append([("•  ", 12.5, True, YELLOW), (it, 12.5, False, WHITE)])
    txt(s, x + 0.3, 2.75, cw - 0.6, 2.8, paras, sp_after=10, line_sp=1.08)
txt(s, 0.7, 6.15, 9.5, 0.9, [[("Électrifier sans déforester", 19, True, YELLOW, False, HEAD),
    (" : une politique d'équité, de santé et de climat à la fois.", 16, False, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE, line_sp=1.05)
txt(s, 0.72, 6.95, 11.9, 0.4, [[("Bastou OURO-TAGBA  ·  Data Challenge Environnement · Défi 2  ·  Togo AI Lab / MESPTN", 11, False, RGBColor(0x9F,0xC4,0xAF))]])

out = os.path.join(BASE, "rapport", "PANORAMA_Rapport.pptx")
prs.save(out)
print("OK -", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
